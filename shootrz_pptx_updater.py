"""
SHOOTRZ Presentation Automation Script
======================================
Production-safe python-pptx script.
ONLY replaces values verified from actual codebase analysis.
Does NOT touch: layout, fonts, colors, positioning.

Author: Architecture Audit — April 2026
"""

import logging
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt

# ── Logging ──────────────────────────────────────────────────────────────────
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s  %(levelname)-8s  %(message)s",
)
log = logging.getLogger("SHOOTRZ-pptx-updater")

# ── Verified Replacement Map ──────────────────────────────────────────────────
# SOURCE OF TRUTH for each value:
#
# VERIFIED (from actual code/runs):
#   - Elbow weight "30%" → real release_mechanics weight = 0.35 (35%)
#     Source: backend/mvp/core/metrics.py line 248 + config component_weights
#   - Knee weight "25%" → real loading_quality weight = 0.30 (30%)
#     Source: backend/mvp/core/metrics.py line 248
#   - Follow-through weight "25%" → real follow_through_control weight = 0.20 (20%)
#     Source: backend/mvp/core/metrics.py line 248
#   - Balance weight "20%" → real balance_stability weight = 0.15 (15%)
#     Source: backend/mvp/core/metrics.py line 248
#   - Pose accuracy "91.4%" → real overall confidence from run 3f4b651f = 84.1%
#     Source: backend/outputs/3f4b651f.../confidence_summary.json
#   - Processing time "4.8s avg" → NOT MEASURED in code; marked as UNVERIFIED → SKIPPED
#   - Score "82", "81/100" → UNVERIFIED (no run produced these) → SKIPPED
#   - Component scores (74, 88, 79) → UNVERIFIED → SKIPPED
#
# UNVERIFIED (demo/synthetic values) — NOT REPLACED:
#   - All performance latency metrics (4.8s, 1.1s, 2.3s, 9.1s, 16.2s)
#   - Concurrency metrics (50/100 users error rates)
#   - Low-light accuracy (61.2%)
#   - Memory usage (218 MB)
#   - Test pass/fail breakdown by module
#   - Overall score 81 / component scores 82, 74, 88, 79
#
VERIFIED_REPLACEMENTS = [
    # ── Slide 6: AI Engine — Score Breakdown weights ──────────────────────────
    # The presentation labels "Elbow Release Angle (30%)" but code has
    # release_mechanics weight = 0.35. Similarly all 4 weights differ.
    {
        "old": "Elbow Release Angle (30%)",
        "new": "Elbow Release Angle (35%)",
        "source": "mvp/core/metrics.py — component_weights.release_mechanics = 0.35",
        "confidence": "HIGH",
    },
    {
        "old": "Knee Bend Depth (25%)",
        "new": "Knee Bend Depth (30%)",
        "source": "mvp/core/metrics.py — component_weights.loading_quality = 0.30",
        "confidence": "HIGH",
    },
    {
        "old": "Follow-Through Arc (25%)",
        "new": "Follow-Through Arc (20%)",
        "source": "mvp/core/metrics.py — component_weights.follow_through_control = 0.20",
        "confidence": "HIGH",
    },
    {
        "old": "Body Alignment (20%)",
        "new": "Body Alignment / Balance (15%)",
        "source": "mvp/core/metrics.py — component_weights.balance_stability = 0.15",
        "confidence": "HIGH",
    },
    # ── Slide 6: Pose Accuracy ────────────────────────────────────────────────
    # 91.4% is NOT a real measured value. Real value from run 3f4b651f = 84.1%.
    # NOTE: This is from 1 run only — treat as indicative, not statistical average.
    {
        "old": "91.4%",
        "new": "84.1%*",
        "source": "outputs/3f4b651f.../confidence_summary.json — overall=0.8413 (1 run sample)",
        "confidence": "MEDIUM — single run, good-lighting conditions",
    },
]

# ── Footnote to add on the slide (informational only — logged, not written) ──
FOOTNOTES_TO_LOG = [
    "* 84.1% is the mean joint-detection confidence from a single 13.7s run (411 frames).",
    "  Processing time, concurrency, low-light, and memory metrics are NOT instrumented",
    "  in the codebase and should be measured before the final presentation.",
    "  Score values (81/100, component scores) are UNVERIFIED — not from test output.",
]


def _replace_in_run(run, old: str, new: str) -> bool:
    """Replace old with new in a single run. Returns True if replacement made."""
    if old in run.text:
        run.text = run.text.replace(old, new)
        return True
    return False


def apply_replacements(pptx_path: Path, out_path: Path) -> dict:
    """
    Load presentation, apply verified replacements, save to new file.
    Returns a report dict with counts.
    """
    prs = Presentation(str(pptx_path))
    total_replacements = 0
    total_skipped = 0
    report = {}

    for replacement in VERIFIED_REPLACEMENTS:
        old = replacement["old"]
        new = replacement["new"]
        source = replacement["source"]
        found_count = 0

        for slide_idx, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if _replace_in_run(run, old, new):
                            found_count += 1
                            log.info(
                                f"  ✓ Slide {slide_idx} | Shape '{shape.name}' | "
                                f"'{old}' → '{new}'"
                            )

        if found_count > 0:
            total_replacements += found_count
            report[old] = {
                "status": "REPLACED",
                "count": found_count,
                "new_value": new,
                "source": source,
            }
            log.info(f"Replaced '{old}' → '{new}' ({found_count}x) | {source}")
        else:
            total_skipped += 1
            report[old] = {
                "status": "NOT_FOUND",
                "count": 0,
                "new_value": new,
                "source": source,
            }
            log.warning(f"  ✗ NOT FOUND: '{old}' (may already be updated or text differs)")

    prs.save(str(out_path))
    log.info(f"\nSaved: {out_path}")
    log.info(f"Summary: {total_replacements} replacements made, {total_skipped} not found.")

    log.info("\n── IMPORTANT NOTES (not written to PPTX) ──────────────────────────────")
    for note in FOOTNOTES_TO_LOG:
        log.info(note)

    return report


def main():
    input_path = Path("SHOOTRZ_PreFinal_Presentation.pptx")
    output_path = Path("SHOOTRZ_PreFinal_Presentation_ActualData.pptx")

    if not input_path.exists():
        log.error(f"Input file not found: {input_path}")
        log.error("Run this script from the same directory as the PPTX file.")
        return

    log.info("=== SHOOTRZ Presentation Updater (Production-Safe) ===")
    log.info(f"Input:  {input_path}")
    log.info(f"Output: {output_path}")
    log.info(f"Replacements planned: {len(VERIFIED_REPLACEMENTS)}")
    log.info("NOTE: Only VERIFIED values are changed. Unverified metrics are left as-is.\n")

    report = apply_replacements(input_path, output_path)

    print("\n── Replacement Report ─────────────────────────────────────────────────────")
    print(f"{'OLD VALUE':<40} {'STATUS':<12} {'NEW VALUE'}")
    print("─" * 90)
    for old_val, info in report.items():
        status = info["status"]
        new_val = info["new_value"]
        print(f"{old_val:<40} {status:<12} {new_val}")
    print("─" * 90)
    print(f"\nOutput saved to: {output_path}")
    print("\nMETRICS THAT REMAIN UNVERIFIED (NOT changed — need instrumentation):")
    unverified = [
        "4.8s avg processing time → no time.perf_counter() in pipeline",
        "1.1s dashboard refresh  → not measured",
        "2.3s cold launch        → not measured",
        "9.1s / 0.4% err @ 50 users  → no load test script exists",
        "16.2s / 8.3% err @ 100 users → no load test script exists",
        "61.2% low-light accuracy → not measured",
        "218 MB peak memory      → no psutil instrumentation",
        "81/100 weighted score   → no real test video produced this",
        "Component scores 82/74/88/79 → unverified demo values",
        "Test module pass rates (Auth/Video/etc) → manual tracking, not from pytest",
    ]
    for item in unverified:
        print(f"  ⚠  {item}")


if __name__ == "__main__":
    main()
